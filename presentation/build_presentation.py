from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTPUT = ROOT / "Wildroot Village - Ecosystem Game Presentation.pptx"

SLIDE_W = 13.333
SLIDE_H = 7.5

BG = "071F18"
BG_ALT = "0A281F"
PANEL = "10372A"
PANEL_ALT = "173F31"
BORDER = "2B5B45"
CREAM = "F3F1D5"
WHITE = "F7FAF4"
MUTED = "B6C9BB"
DIM = "7F9B89"
GREEN = "A7E87D"
GREEN_DARK = "4F9366"
GOLD = "E4C967"
CORAL = "F47E70"
BLUE = "6FC3D5"
PURPLE = "C79AEC"

TITLE_FONT = "Georgia"
BODY_FONT = "Aptos"


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_bg(slide, color: str = BG) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 1.0,
    radius: bool = False,
):
    if radius:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(fill)
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = rgb(line)
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def style_run(run, size: float, color: str, bold: bool = False, font: str = BODY_FONT, italic: bool = False) -> None:
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 18,
    color: str = WHITE,
    bold: bool = False,
    font: str = BODY_FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0.02,
    italic: bool = False,
    fit: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    paragraph = tf.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    style_run(run, size, color, bold, font, italic)
    if fit:
        tf.fit_text(max_size=Pt(size))
    return box


def add_rich_lines(slide, lines, x, y, w, h, margin=0.08, valign=MSO_ANCHOR.TOP):
    """lines: list of dicts with runs=[(text, size, color, bold, font, italic)]."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    tf.word_wrap = True
    for index, spec in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        p.space_before = Pt(spec.get("before", 0))
        p.space_after = Pt(spec.get("after", 4))
        p.line_spacing = spec.get("line_spacing", 1.0)
        for run_spec in spec["runs"]:
            text_value, size, color, bold, font, italic = run_spec
            run = p.add_run()
            run.text = text_value
            style_run(run, size, color, bold, font, italic)
    return box


def add_label(slide, text: str, x: float, y: float, w: float, color: str = GREEN) -> None:
    add_text(slide, text.upper(), x, y, w, 0.26, 9.5, color, True, BODY_FONT)


def add_title(slide, title: str, kicker: str, number: int, subtitle: str | None = None) -> None:
    add_label(slide, kicker, 0.52, 0.28, 5.8)
    add_text(slide, title, 0.52, 0.54, 11.7, 0.52, 28, CREAM, True, TITLE_FONT)
    if subtitle:
        add_text(slide, subtitle, 0.54, 1.02, 11.6, 0.32, 11.2, MUTED)
    add_text(slide, f"0{number}", 12.32, 0.37, 0.48, 0.30, 10, GREEN, True, BODY_FONT, PP_ALIGN.RIGHT)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 12.91, 0.30, 0.03, 0.55, GREEN)


def add_footer(slide, number: int, text_value: str = "WILDROOT VILLAGE · ECOSYSTEM VILLAGE PROJECT") -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.52, 7.17, 12.29, 0.012, BORDER)
    add_text(slide, text_value, 0.52, 7.21, 7.2, 0.18, 7.5, DIM, True)
    add_text(slide, str(number), 12.45, 7.19, 0.35, 0.18, 8.5, DIM, True, align=PP_ALIGN.RIGHT)


def add_card(slide, x, y, w, h, fill=PANEL, line=BORDER, radius=True):
    return add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, 0.8, radius)


def add_picture_cover(slide, path: Path, x, y, w, h, border=True):
    with Image.open(path) as image:
        image_w, image_h = image.size
    image_ratio = image_w / image_h
    box_ratio = w / h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    if image_ratio > box_ratio:
        visible = box_ratio / image_ratio
        crop = (1 - visible) / 2
        picture.crop_left = crop
        picture.crop_right = crop
    else:
        visible = image_ratio / box_ratio
        crop = (1 - visible) / 2
        picture.crop_top = crop
        picture.crop_bottom = crop
    if border:
        outline = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, None, BORDER, 1.3)
        return picture, outline
    return picture


def add_picture_region(slide, path: Path, x, y, w, h, region, border=True):
    """Place a source-image pixel region into a destination box."""
    with Image.open(path) as image:
        image_w, image_h = image.size
    left, top, right, bottom = region
    crop_left = left / image_w
    crop_right = (image_w - right) / image_w
    crop_top = top / image_h
    crop_bottom = (image_h - bottom) / image_h
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    picture.crop_left = crop_left
    picture.crop_right = crop_right
    picture.crop_top = crop_top
    picture.crop_bottom = crop_bottom
    if border:
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, None, BORDER, 1.3)
    return picture


def add_number_marker(slide, number: int, x: float, y: float, color: str = GREEN) -> None:
    add_shape(slide, MSO_SHAPE.OVAL, x, y, 0.33, 0.33, color, BG, 1.2)
    add_text(slide, str(number), x, y + 0.015, 0.33, 0.24, 10.2, BG, True, BODY_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)


def set_notes(slide, notes: str) -> None:
    # Keep notes in the companion markdown file. python-pptx-generated notes
    # currently make macOS Keynote/Quick Look stall while importing the deck.
    # Leaving the PowerPoint notes layer untouched keeps the presentation
    # compatible with both Keynote and PowerPoint.
    return None


def slide_1(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    add_shape(slide, MSO_SHAPE.OVAL, -1.35, 5.58, 4.0, 2.7, "0D3024")
    add_shape(slide, MSO_SHAPE.OVAL, 2.05, -0.90, 2.0, 1.25, "0D3024")
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.58, 0.65, 0.05, 0.82, GREEN)
    add_label(slide, "Ecosystem Village project", 0.83, 0.68, 4.5)
    add_text(slide, "Wildroot\nVillage", 0.80, 1.18, 4.8, 1.65, 43, CREAM, True, TITLE_FONT)
    add_text(slide, "An ecosystem-management and environmental learning game", 0.82, 3.04, 4.62, 0.70, 18.5, WHITE, False, BODY_FONT)
    add_text(slide, "Grow a home in the forest.\nLeave enough wildness for tomorrow.", 0.82, 3.92, 4.82, 0.78, 17.5, GREEN, False, TITLE_FONT, italic=True)

    chips = [("BUILD", GREEN), ("SIMULATE", BLUE), ("LEARN", GOLD)]
    x_pos = 0.82
    for label, color in chips:
        add_card(slide, x_pos, 5.05, 1.26, 0.42, PANEL_ALT, color)
        add_text(slide, label, x_pos, 5.13, 1.26, 0.16, 8.8, color, True, BODY_FONT, PP_ALIGN.CENTER)
        x_pos += 1.42

    add_card(slide, 0.82, 5.60, 1.62, 0.34, BG_ALT, GREEN)
    add_text(slide, "≈ 3-MINUTE PRESENTATION", 0.82, 5.69, 1.62, 0.11, 6.9, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    add_text(slide, "CREATED BY", 0.82, 6.16, 2.6, 0.18, 8.5, DIM, True)
    add_text(slide, "James · age 13", 0.82, 6.37, 2.8, 0.23, 11.2, CREAM, True)
    add_text(slide, "PRESENTED BY", 0.82, 6.65, 2.6, 0.16, 7.8, DIM, True)
    add_text(slide, "Daniel Bryar", 0.82, 6.83, 2.8, 0.20, 9.8, GREEN, True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.94, 0.43, 6.87, 6.60, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "title-screen.png", 6.05, 0.54, 6.66, 6.38)
    add_card(slide, 10.82, 6.37, 1.57, 0.37, BG_ALT, GREEN)
    add_text(slide, "OPENING SCREEN", 10.82, 6.46, 1.57, 0.13, 7.6, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    set_notes(slide, """
Introduce Wildroot Village as the game created in the Ecosystem Village project. It is a browser-based settlement simulation with an environmental learning purpose. The central promise is captured by the line “Grow a home in the forest; leave enough wildness for tomorrow.” The player is not only trying to make a large village. They must provide food, water, housing and materials while protecting the living system that makes survival possible. The presentation uses screenshots from the working game and explains the simulation, controls, ecology, scenarios and learning design. The game can be opened directly from index.html, needs no account or installation, and saves the current village, learning progress and achievements locally in the browser.
""")
    return slide


def slide_creator_story(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Why I created Wildroot Village", "Creator story · James, age 13", 2,
              "An idea inspired by climate change, global warming and the power of environmental education.")

    add_card(slide, 0.52, 1.46, 5.15, 4.68, PANEL, BORDER)
    add_card(slide, 0.78, 1.72, 2.05, 0.40, BG_ALT, GREEN)
    add_text(slide, "JAMES  ·  CREATOR  ·  AGE 13", 0.78, 1.82, 2.05, 0.13, 7.8, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)
    add_card(slide, 2.98, 1.72, 2.38, 0.40, BG_ALT, GOLD)
    add_text(slide, "PRESENTED BY  ·  DANIEL BRYAR", 2.98, 1.82, 2.38, 0.13, 7.1, GOLD, True, BODY_FONT, PP_ALIGN.CENTER)
    add_text(slide,
             "“Yesterday, thinking about climate change and global warming gave me an idea: a game that helps kids see how every building choice can change the environment.”",
             0.82, 2.30, 4.54, 1.40, 20.5, CREAM, False, TITLE_FONT, italic=True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.82, 3.92, 0.48, 0.04, GREEN)
    add_text(slide,
             "I want children to learn by making decisions and watching the results—not only by reading facts. Learning is integrated into the game through the Eco Coach, field missions, knowledge checks, inspections and scenarios.",
             0.82, 4.13, 4.54, 1.02, 11.8, MUTED)
    add_card(slide, 0.82, 5.31, 4.54, 0.58, BG_ALT, GREEN)
    add_text(slide, "OVERALL GOAL", 1.02, 5.43, 1.10, 0.15, 7.7, GREEN, True)
    add_text(slide, "Create a thriving village with a healthy ecosystem.", 2.12, 5.39, 3.03, 0.24, 10.2, CREAM, True)

    add_card(slide, 5.92, 1.46, 6.89, 4.68, BG_ALT, BORDER)
    add_label(slide, "One simple building can change a whole system", 6.22, 1.72, 5.95)

    # Central building illustration.
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, 8.85, 2.73, 1.30, 0.82, GOLD, GOLD)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.00, 3.26, 1.00, 0.92, "D9D7B9", GOLD, 1.0)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.39, 3.64, 0.23, 0.54, PANEL_ALT, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.12, 3.44, 0.18, 0.18, BLUE, None)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 9.69, 3.44, 0.18, 0.18, BLUE, None)
    add_text(slide, "A SIMPLE\nBUILDING", 8.77, 4.28, 1.46, 0.46, 9.4, GOLD, True, BODY_FONT, PP_ALIGN.CENTER)

    impacts = [
        ("HABITAT", "clear it or protect it", GREEN, 6.25, 2.26),
        ("WATER", "consume, capture or pollute", BLUE, 10.55, 2.26),
        ("AIR + NOISE", "create emissions or disturbance", PURPLE, 6.25, 3.72),
        ("SOIL", "compact, erode or restore", CORAL, 10.55, 3.72),
        ("PEOPLE", "change jobs, health and learning", GOLD, 8.23, 5.13),
    ]
    for heading, body, color, x, y in impacts:
        width = 2.25 if heading != "PEOPLE" else 2.78
        add_card(slide, x, y, width, 0.72, PANEL, color)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.14, y + 0.17, 0.22, 0.22, color)
        add_text(slide, heading, x + 0.47, y + 0.12, width - 0.60, 0.17, 8.6, color, True)
        add_text(slide, body, x + 0.47, y + 0.35, width - 0.60, 0.20, 7.7, MUTED)

    add_text(slide, "→", 8.52, 2.45, 0.28, 0.25, 16, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)
    add_text(slide, "←", 10.23, 2.45, 0.28, 0.25, 16, BLUE, True, BODY_FONT, PP_ALIGN.CENTER)
    add_text(slide, "→", 8.52, 3.90, 0.28, 0.25, 16, PURPLE, True, BODY_FONT, PP_ALIGN.CENTER)
    add_text(slide, "←", 10.23, 3.90, 0.28, 0.25, 16, CORAL, True, BODY_FONT, PP_ALIGN.CENTER)
    add_text(slide, "↑", 9.36, 4.82, 0.28, 0.25, 16, GOLD, True, BODY_FONT, PP_ALIGN.CENTER)

    add_card(slide, 0.52, 6.38, 12.29, 0.47, BG_ALT, BORDER)
    add_text(slide, "LEARNING IS INTEGRATED INTO PLAY", 0.76, 6.49, 2.48, 0.15, 8.1, GREEN, True)
    add_text(slide, "Eco Coach   ·   field missions   ·   knowledge checks   ·   building inspections   ·   environmental scenarios",
             3.24, 6.48, 9.27, 0.16, 8.6, CREAM, True, BODY_FONT, PP_ALIGN.CENTER)

    add_footer(slide, 2, "CREATED BY JAMES, AGE 13 · PRESENTED BY DANIEL BRYAR")
    set_notes(slide, """
I came up with Wildroot Village yesterday while thinking about climate change and global warming. I wanted to create something that could help children understand the environment through choices and consequences. A single building may look simple, but it can remove or protect habitat, use or capture water, compact or restore soil, create air pollution or noise, and affect people’s health and learning. This is why education is integrated into the game rather than placed in a separate textbook section. The Eco Coach, field missions, knowledge checks, inspections and scenarios explain what is happening while the player builds. The overall goal is to prove that a village does not need to choose between people and nature: the best result is a thriving community supported by a healthy ecosystem.
""")
    return slide


def slide_2(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "The idea: prosper without breaking the system", "Core game loop", 3,
              "Every expansion solves one problem while creating new pressure somewhere else.")

    steps = [
        ("01", "PLAN", "Read the map, storage, population and weakest ecosystem indicator.", GREEN),
        ("02", "BUILD", "Place homes, services, workplaces and nature-based solutions.", BLUE),
        ("03", "MONITOR", "Watch 24-hour resource trends, staff, noise, pollution and weather.", GOLD),
        ("04", "ADAPT", "Change layout, jobs, reserves and restoration before a threshold is crossed.", PURPLE),
    ]
    card_w = 2.86
    x_positions = [0.52, 3.63, 6.74, 9.85]
    for index, ((num, label, body, color), x) in enumerate(zip(steps, x_positions)):
        add_card(slide, x, 1.48, card_w, 1.54, PANEL, BORDER)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.18, 1.67, 0.44, 0.44, color)
        add_text(slide, num, x + 0.18, 1.76, 0.44, 0.16, 9.0, BG, True, BODY_FONT, PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.76, 1.64, 1.76, 0.26, 13.5, color, True)
        add_text(slide, body, x + 0.20, 2.19, card_w - 0.40, 0.62, 10.7, MUTED)
        if index < 3:
            add_text(slide, "→", x + 2.85, 1.97, 0.28, 0.32, 20, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    outcomes = [
        ("THE GOAL", "Grow a functioning village", "Balance population, food, water, timber, stone, health, morale and education—without sacrificing the forest.", GREEN),
        ("THE MAIN TENSION", "Short-term output vs. long-term resilience", "Industry can speed construction, but resource use, habitat loss, pollution and noise can weaken future production.", GOLD),
        ("THE FAILURE RULE", "A weak link can collapse the settlement", "The village is lost below 5% overall ecosystem health, or immediately when any one of the six indicators reaches 0%.", CORAL),
    ]
    x_positions = [0.52, 4.67, 8.82]
    for (label, heading, body, color), x in zip(outcomes, x_positions):
        add_card(slide, x, 3.42, 3.97, 2.32, PANEL_ALT, color)
        add_label(slide, label, x + 0.22, 3.64, 3.35, color)
        add_text(slide, heading, x + 0.22, 3.98, 3.52, 0.51, 16.5, CREAM, True, TITLE_FONT)
        add_text(slide, body, x + 0.22, 4.60, 3.52, 0.82, 11.3, MUTED)

    add_card(slide, 0.52, 6.03, 12.29, 0.73, BG_ALT, BORDER)
    facts = [
        ("100 × 100", "forest world"),
        ("3 minutes", "per day at 1×"),
        ("1× / 2× / 3×", "true simulation speed"),
        ("Freeplay + 8", "challenge scenarios"),
        ("Local", "browser save"),
    ]
    x = 0.76
    for idx, (value, label) in enumerate(facts):
        add_text(slide, value, x, 6.17, 1.62, 0.22, 11.6, GREEN, True)
        add_text(slide, label, x, 6.42, 1.75, 0.16, 8.3, DIM)
        if idx < len(facts) - 1:
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + 1.96, 6.16, 0.012, 0.39, BORDER)
        x += 2.43

    add_footer(slide, 3)
    set_notes(slide, """
Explain the loop from left to right. Planning begins with evidence: resource totals, daily trends, workforce availability, weather and the six ecosystem indicators. Building then converts plans into a physical settlement. The player monitors what actually happens over a full day, because staffed workplaces run during the day while automatic buildings, consumption, weather and natural processes continue overnight. Finally, the player adapts by changing placement, workforce, storage, production or restoration. Success means a village that can keep growing without destroying its own support system. The important design choice is the weak-link rule: a high average cannot hide a totally collapsed river, soil system or habitat. The game ends if the average ecosystem drops below 5%, or if forest cover, wildlife, water quality, soil health, clean air or biodiversity reaches zero.
""")
    return slide


def slide_3(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "How the player controls the village", "Interface tour", 4,
              "The screen keeps construction, simulation evidence and environmental feedback visible at the same time.")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.42, 1.39, 8.86, 5.03, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "village-gameplay.png", 0.51, 1.48, 8.68, 4.86)
    marker_positions = [(1, 3.44, 1.61), (2, 0.72, 2.83), (3, 4.81, 3.42), (4, 8.65, 2.73), (5, 5.74, 5.83)]
    for number, x, y in marker_positions:
        add_number_marker(slide, number, x, y)

    legend = [
        (1, "Resources + time", "Population, stored materials, /day trends, season, weather and speed controls."),
        (2, "Planning desk", "Choose Village, Work or Nature buildings; then Inspect or Destroy."),
        (3, "Living map", "Drag and zoom across the 100 × 100 forest; place only on free clearing tiles."),
        (4, "Ecosystem + Eco Coach", "Six indicators, overall score, weakest trend, major pressure and best support."),
        (5, "Workforce + footprint", "Jobs, available workers, children, schooling, health, morale and settlement pressure."),
    ]
    y = 1.47
    for number, heading, body in legend:
        add_card(slide, 9.49, y, 3.32, 0.82, PANEL, BORDER)
        add_number_marker(slide, number, 9.66, y + 0.20)
        add_text(slide, heading, 10.10, y + 0.13, 2.43, 0.22, 11.2, CREAM, True)
        add_text(slide, body, 10.10, y + 0.39, 2.48, 0.34, 8.8, MUTED)
        y += 0.94

    add_card(slide, 9.49, 6.15, 3.32, 0.55, BG_ALT, GREEN)
    add_text(slide, "DRAG pan  ·  SCROLL zoom  ·  O/P rotate", 9.65, 6.28, 2.98, 0.14, 7.8, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)
    add_text(slide, "SPACE pause  ·  ESC cancel  ·  HOLD tree priority", 9.65, 6.47, 2.98, 0.12, 7.2, MUTED, False, BODY_FONT, PP_ALIGN.CENTER)

    add_footer(slide, 4)
    set_notes(slide, """
Walk around the screen clockwise. The header shows the current day, season, time, weather and simulation speed. The resource strip reports people, food, water, timber, stone and ecosystem health, including weighted daily trends. The Planning Desk on the left contains build categories and the Inspect and Destroy tools. The central canvas is a zoomable 100 by 100 forest world with an irregular clearing; players drag to pan, scroll to zoom, use O and P to rotate, click residents to inspect them, and can hold a tree while Inspect is active to make it a logging priority. The right panel explains the current ecosystem and Steward’s Path goals. The footer reports jobs, available workers, children, education, health, morale and settlement footprint. Space pauses or resumes and Escape cancels the current tool. Descriptions can be turned off, and the preference is saved.
""")
    return slide


def slide_4(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "A living model, not a single score", "Ecosystem simulation", 5,
              "Six connected indicators reveal weak links that an average could otherwise hide.")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.42, 1.39, 7.04, 4.05, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "field-guide.png", 0.51, 1.48, 6.86, 3.87)
    add_card(slide, 0.73, 4.90, 2.20, 0.29, BG_ALT, GREEN)
    add_text(slide, "LIVE 24-HOUR FIELD GUIDE", 0.73, 4.97, 2.20, 0.11, 7.2, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    add_label(slide, "The six indicators", 7.72, 1.42, 2.7)
    indicators = [
        ("Forest cover", "Habitat, carbon and recovery capacity", GREEN, 76),
        ("Wildlife", "Abundance and security of animal populations", GOLD, 71),
        ("Water quality", "Cleanliness, supply and watershed health", BLUE, 67),
        ("Soil health", "Fertility, soil life and structural stability", CORAL, 65),
        ("Clean air", "Smoke, dust and atmospheric capacity", PURPLE, 72),
        ("Biodiversity", "Species, habitats and ecological relationships", GREEN, 70),
    ]
    for index, (name, body, color, score) in enumerate(indicators):
        col = index % 2
        row = index // 2
        x = 7.72 + col * 2.56
        y = 1.77 + row * 1.02
        add_card(slide, x, y, 2.38, 0.88, PANEL, BORDER)
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.13, y + 0.14, 0.23, 0.23, color)
        add_text(slide, name, x + 0.46, y + 0.10, 1.43, 0.20, 10.6, CREAM, True)
        add_text(slide, f"{score}", x + 1.90, y + 0.10, 0.31, 0.18, 10.3, color, True, BODY_FONT, PP_ALIGN.RIGHT)
        add_text(slide, body, x + 0.16, y + 0.40, 2.06, 0.31, 8.2, MUTED)

    add_card(slide, 0.52, 5.70, 12.29, 1.11, BG_ALT, BORDER)
    add_label(slide, "Cause and effect", 0.73, 5.86, 1.65)
    flow = [
        ("BUILDINGS + PEOPLE", "use resources; create output, noise and pollution", GREEN),
        ("WEATHER + SEASONS", "change demand, production, fires, floods and growth", BLUE),
        ("24-HOUR RATES", "combine daytime staffing with overnight processes", GOLD),
        ("FEEDBACK", "changes food, health, morale and long-term resilience", PURPLE),
    ]
    x = 2.20
    for index, (heading, body, color) in enumerate(flow):
        add_text(slide, heading, x, 5.84, 2.15, 0.18, 8.4, color, True)
        add_text(slide, body, x, 6.08, 2.20, 0.35, 8.1, MUTED)
        if index < len(flow) - 1:
            add_text(slide, "→", x + 2.18, 6.05, 0.24, 0.23, 15, DIM, True, BODY_FONT, PP_ALIGN.CENTER)
        x += 2.51
    add_card(slide, 10.73, 6.38, 1.81, 0.25, "38251F", CORAL)
    add_text(slide, "COLLAPSE: <5% OR ANY 0%", 10.73, 6.44, 1.81, 0.10, 6.6, CORAL, True, BODY_FONT, PP_ALIGN.CENTER)

    add_footer(slide, 5)
    set_notes(slide, """
The ecosystem is represented by forest cover, wildlife, water quality, soil health, clean air and biodiversity. These are deliberately separate because an acceptable average can conceal one failing system. The Eco Coach reads the real values and full 24-hour rates, identifies the weakest or fastest-falling indicator, names the strongest modelled pressure and recovery source, and connects it to other parts of the system. Rates are weighted across staffed daytime work and automatic overnight processes, so the advice does not reverse just because night begins. Buildings, population, weather, seasons, habitat, pollution and restoration all feed the model. The Field Guide also states an important limitation: the six scores are simplified indicators. Real environmental decisions require field measurements, local and Indigenous knowledge, uncertainty and impacts beyond the visible map.
""")
    return slide


def slide_5(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "People, work and the resource economy", "Village systems", 6,
              "Residents have homes, ages, jobs and journeys; buildings operate on different schedules and capacities.")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.42, 1.39, 6.13, 4.78, "061A14", None, radius=True)
    add_picture_region(slide, ROOT / "village-gameplay.png", 0.51, 1.48, 5.95, 4.60, (245, 130, 1320, 900))
    add_card(slide, 0.73, 5.60, 1.58, 0.29, BG_ALT, GREEN)
    add_text(slide, "A WORKING VILLAGE", 0.73, 5.67, 1.58, 0.11, 7.1, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    cards = [
        ("PEOPLE + LIFE", GREEN,
         "Freeplay starts with 10 adult founders. The population includes children, adults and elders. Children become adults after 7 in-game days; residents have persistent homes and a hidden natural lifespan of 40–60 days."),
        ("WORK + MOVEMENT", BLUE,
         "Staffed buildings operate 07:00–18:30. Residents walk between home, work, parks, school and storage; workers carry 1–3 items per supply trip. Automatic buildings continue overnight."),
        ("STORAGE + SUPPLY", GOLD,
         "The Founders’ Hearth stores 400 of each material. Every Storehouse adds 200 food, water, timber and stone capacity. Production above the cap is discarded, so overnight reserves matter."),
    ]
    y = 1.48
    for heading, color, body in cards:
        add_card(slide, 6.80, y, 6.01, 1.30, PANEL, BORDER)
        add_shape(slide, MSO_SHAPE.RECTANGLE, 6.80, y, 0.06, 1.30, color)
        add_text(slide, heading, 7.05, y + 0.15, 2.17, 0.22, 10.0, color, True)
        add_text(slide, body, 7.05, y + 0.45, 5.48, 0.67, 10.1, MUTED)
        y += 1.45

    add_label(slide, "Three important trade-offs", 6.82, 5.88, 3.2)
    tradeoffs = [
        ("FARMS", "1 farmer = 0.5×; 2 = 1×; 3 = 2× output—and 2× water use and pressure.", GREEN),
        ("FORESTRY", "Logging is fast nearby; Wood Farms regrow 16 plots in 5 days. Full storage pauses cutting.", GOLD),
        ("SITING", "Industry near crops reduces yield; loud buildings near occupied homes lower morale and health.", PURPLE),
    ]
    x = 6.80
    for heading, body, color in tradeoffs:
        add_card(slide, x, 6.18, 1.90, 0.65, BG_ALT, color)
        add_text(slide, heading, x + 0.12, 6.29, 1.62, 0.15, 7.5, color, True)
        add_text(slide, body, x + 0.12, 6.47, 1.65, 0.28, 6.9, MUTED)
        x += 2.05

    add_footer(slide, 6)
    set_notes(slide, """
Residents are simulated individuals rather than a single population number. They keep persistent names, ages, housing and workplace assignments. Children mature after seven in-game days; schools can educate up to sixteen children and improve long-term productivity and environmental practice. At night, residents enter their assigned homes and become hidden until after 07:00. Staffed workplaces operate from 07:00 to 18:30, while housing, wells, storehouses, parks, sanctuaries, windmills and other automatic services can continue overnight. Workers carry one to three items and make visible supply trips. Storage is a real constraint: the Hearth holds 400 of each main resource and each Storehouse adds 200. Farming, forestry and settlement layout all create trade-offs. Extra farm labour doubles output but also water use and ecological pressure; Wood Farms create renewable timber; and pollution or noise has stronger effects when buildings are badly placed.
""")
    return slide


def slide_6(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Eight scenarios turn systems into problems", "Scenarios, weather and risk", 7,
              "Each village begins with a different crisis, four measurable goals and a named environmental learning focus.")

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.42, 1.39, 5.95, 3.39, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "scenario-list.png", 0.51, 1.48, 5.77, 3.21)
    add_card(slide, 0.73, 4.28, 1.38, 0.27, BG_ALT, GREEN)
    add_text(slide, "8 STARTING VILLAGES", 0.73, 4.35, 1.38, 0.10, 6.8, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 6.52, 1.39, 6.29, 3.39, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "restoration-briefing.png", 6.61, 1.48, 6.11, 3.21)
    add_card(slide, 10.69, 4.28, 1.75, 0.27, "38251F", CORAL)
    add_text(slide, "RESTORATION BRIEFING", 10.69, 4.35, 1.75, 0.10, 6.6, CORAL, True, BODY_FONT, PP_ALIGN.CENTER)

    columns = [
        ("FIVE CHALLENGE VILLAGES", GREEN,
         "Winter Watch · The Dry River · Timber Debt · Smoke Valley · The Green City\n\nThey test seasonal reserves, drought, renewable forestry, clean industry and dense sustainable growth."),
        ("THREE ACTIVE EMERGENCIES", CORAL,
         "After the Fire · The Poisoned River · The Silent Fields\n\nEach starts below 40% and loses 1.0 ecosystem point/day until its containment actions are completed."),
        ("DYNAMIC RISK", BLUE,
         "Weather ranges from ~2-hour thunderstorms and ~12-hour rain to 3-day heatwaves, 4-day cloud systems and 7-day droughts. Random events arrive every ⅓–7 days."),
    ]
    x = 0.52
    for heading, color, body in columns:
        add_card(slide, x, 5.06, 3.96, 1.70, PANEL, BORDER)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, 5.06, 3.96, 0.06, color)
        add_text(slide, heading, x + 0.20, 5.27, 3.50, 0.22, 9.1, color, True)
        add_text(slide, body, x + 0.20, 5.59, 3.50, 0.93, 9.3, MUTED)
        x += 4.16

    add_footer(slide, 7)
    set_notes(slide, """
The scenario selector offers eight prebuilt villages. Winter Watch focuses on preparing seasonal reserves; The Dry River on water capture and quality; Timber Debt on replacing destructive logging; Smoke Valley on separating polluters from homes and crops; and The Green City on supporting a dense, educated population with healthy habitat. After the Fire, The Poisoned River and The Silent Fields are restoration emergencies. They begin below 40% ecosystem health and continue losing one overall point per day until specific containment actions are completed. The After the Fire example starts at 34% after a burned watershed. The player must retire every Logging Camp, maintain two Rain Gardens and a Compost Yard, restore forest and soil to at least 65, then survive to Day 30 with at least 60 overall. Seasons, weather and random events keep every solution dynamic. Events include wildfire, flooding, illness, pollinator decline, invasive vines, migration, wildlife conflict, shortages and trade.
""")
    return slide


def slide_7(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_title(slide, "Learning is built into the decisions", "Learning, progression and success", 8,
              "Players learn by observing consequences, testing ideas and explaining what the model leaves out.")

    add_card(slide, 0.52, 1.47, 4.61, 4.63, PANEL, BORDER)
    add_label(slide, "The learning layer", 0.78, 1.71, 2.7)
    learning_points = [
        ("20 field missions", "demonstrate principles through normal gameplay; every lesson links to a real-world example."),
        ("22 knowledge checks", "pause time, explain correct and incorrect answers, and give no resource reward or penalty."),
        ("Learning everywhere", "Eco Coach advice, building inspections, events and scenario briefings all include a learning lens."),
        ("Model literacy", "the game names its simplifications and points to field evidence, local and Indigenous knowledge and uncertainty."),
        ("Saved progression", "learning and achievements persist locally with the current village."),
    ]
    y = 2.07
    colors = [GREEN, GOLD, BLUE, PURPLE, GREEN]
    for index, ((heading, body), color) in enumerate(zip(learning_points, colors), start=1):
        add_shape(slide, MSO_SHAPE.OVAL, 0.78, y + 0.02, 0.29, 0.29, color)
        add_text(slide, str(index), 0.78, y + 0.075, 0.29, 0.10, 7.0, BG, True, BODY_FONT, PP_ALIGN.CENTER)
        add_rich_lines(slide, [{"runs": [
            (heading + " — ", 10.4, CREAM, True, BODY_FONT, False),
            (body, 9.5, MUTED, False, BODY_FONT, False),
        ], "after": 0, "line_spacing": 1.0}], 1.18, y - 0.02, 3.60, 0.54, margin=0.01)
        y += 0.72

    add_card(slide, 0.52, 6.28, 4.61, 0.56, BG_ALT, GREEN)
    add_text(slide, "FINAL CHAPTER", 0.76, 6.39, 1.16, 0.15, 7.8, GREEN, True)
    add_text(slide, "Day 100  ·  nature buildings  ·  every indicator 90+  ·  100% ecosystem", 1.96, 6.38, 2.92, 0.24, 8.0, CREAM, True)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.42, 1.39, 7.39, 3.68, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "field-guide.png", 5.51, 1.48, 7.21, 3.50)
    add_card(slide, 5.72, 4.57, 1.43, 0.27, BG_ALT, GREEN)
    add_text(slide, "FIELD MISSIONS", 5.72, 4.64, 1.43, 0.10, 6.9, GREEN, True, BODY_FONT, PP_ALIGN.CENTER)

    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.42, 4.65, 5.39, 2.10, "061A14", None, radius=True)
    add_picture_cover(slide, ROOT / "knowledge-check.png", 7.51, 4.74, 5.21, 1.92)
    add_card(slide, 10.83, 6.25, 1.52, 0.25, BG_ALT, GOLD)
    add_text(slide, "NO GAMEPLAY PENALTY", 10.83, 6.31, 1.52, 0.10, 6.4, GOLD, True, BODY_FONT, PP_ALIGN.CENTER)

    add_text(slide, "BUILD  →  MEASURE  →  ADAPT  →  LEAVE WILDNESS FOR TOMORROW", 5.50, 5.36, 1.71, 1.09, 13.5, GREEN, True, BODY_FONT, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)

    add_footer(slide, 8, "WILDROOT VILLAGE · BUILD A RESILIENT FUTURE")
    set_notes(slide, """
The learning content is optional but embedded throughout play. Twenty field missions ask the player to demonstrate environmental principles using the normal simulation. Twenty-two knowledge checks pause time, explain both correct and incorrect responses, and never add or remove resources, so learning does not become a shortcut or punishment. Building inspections include an environmental lesson and a decision question; events and scenario briefings also name the concept being tested. Progress continues through five Steward’s Path chapters. The final chapter, A Living Legacy, asks the player to reach Day 100, maintain a Compost Yard, Rain Garden, Wild Sanctuary and Town Park, restore every indicator to at least 90, and achieve a displayed 100% ecosystem score. Achievements persist across villages, and Perfect Balance is unlocked by reaching 100% on Harsh difficulty. End with the game’s core message: build, measure, adapt and leave enough wildness for the future.
""")
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    prs.core_properties.title = "Wildroot Village — Ecosystem Game Presentation"
    prs.core_properties.subject = "A three-minute, eight-slide overview of the Ecosystem Village game"
    prs.core_properties.author = "James, age 13"
    prs.core_properties.keywords = "ecosystem, village, game, sustainability, environmental learning"

    slide_1(prs)
    slide_creator_story(prs)
    slide_2(prs)
    slide_3(prs)
    slide_4(prs)
    slide_5(prs)
    slide_6(prs)
    slide_7(prs)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    result = build()
    print(result)
